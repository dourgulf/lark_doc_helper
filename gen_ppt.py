"""Generate a PowerPoint presentation from output.md"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Color palette
COLOR_BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)       # dark navy
COLOR_BG_ACCENT = RGBColor(0x16, 0x21, 0x3E)      # slightly lighter navy
COLOR_ACCENT = RGBColor(0x0F, 0x3C, 0x78)         # blue accent bar
COLOR_HIGHLIGHT = RGBColor(0x4A, 0xC2, 0xF0)      # cyan highlight
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xCC, 0xD6, 0xE0)
COLOR_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
COLOR_GREEN = RGBColor(0x6C, 0xD9, 0x7E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                 wrap=True):
    if color is None:
        color = COLOR_WHITE
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txBox


def add_bullet_slide(prs, title, bullets, subtitle=None):
    """Add a slide with title + bullet list."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, COLOR_BG_DARK)

    # Left accent bar
    add_rect(slide, 0, 0, 0.08, 7.5, COLOR_HIGHLIGHT)

    # Title background strip
    add_rect(slide, 0.08, 0, 13.25, 1.4, COLOR_BG_ACCENT)

    # Title text
    add_text_box(slide, title, 0.25, 0.15, 12.5, 1.1,
                 font_size=32, bold=True, color=COLOR_HIGHLIGHT)

    if subtitle:
        add_text_box(slide, subtitle, 0.25, 1.5, 12.5, 0.5,
                     font_size=16, color=COLOR_LIGHT_GRAY)

    # Bullets
    y_start = 1.6 if not subtitle else 2.1
    y_spacing = 0.75
    for i, bullet in enumerate(bullets):
        y = y_start + i * y_spacing
        # Bullet marker
        add_text_box(slide, "▶", 0.3, y, 0.4, 0.6,
                     font_size=14, color=COLOR_HIGHLIGHT)
        add_text_box(slide, bullet, 0.75, y, 12.0, 0.65,
                     font_size=18, color=COLOR_WHITE)

    return slide


def add_two_column_slide(prs, title, left_items, right_items,
                         left_header="", right_header=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 0.08, 7.5, COLOR_HIGHLIGHT)
    add_rect(slide, 0.08, 0, 13.25, 1.4, COLOR_BG_ACCENT)
    add_text_box(slide, title, 0.25, 0.15, 12.5, 1.1,
                 font_size=32, bold=True, color=COLOR_HIGHLIGHT)

    # Column divider
    add_rect(slide, 6.7, 1.5, 0.04, 5.7, COLOR_ACCENT)

    col_configs = [
        (0.3, left_header, left_items),
        (6.9, right_header, right_items),
    ]
    for x, header, items in col_configs:
        y = 1.55
        if header:
            add_text_box(slide, header, x, y, 6.0, 0.5,
                         font_size=20, bold=True, color=COLOR_YELLOW)
            y += 0.6
        for item in items:
            add_text_box(slide, f"• {item}", x, y, 6.0, 0.65,
                         font_size=16, color=COLOR_WHITE)
            y += 0.7

    return slide


def build_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: Title ──────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)

    # Top decorative bar
    add_rect(slide, 0, 0, 13.33, 0.12, COLOR_HIGHLIGHT)
    # Bottom decorative bar
    add_rect(slide, 0, 7.38, 13.33, 0.12, COLOR_HIGHLIGHT)
    # Left accent
    add_rect(slide, 0, 0.12, 0.12, 7.26, COLOR_ACCENT)

    # Main title
    add_text_box(slide, "OT Flutter 客户端", 1.0, 1.5, 11.0, 1.2,
                 font_size=52, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "AI 提效工具实践", 1.0, 2.7, 11.0, 1.0,
                 font_size=48, bold=True, color=COLOR_HIGHLIGHT, align=PP_ALIGN.CENTER)

    # Subtitle
    add_rect(slide, 3.5, 4.0, 6.33, 0.06, COLOR_HIGHLIGHT)
    add_text_box(slide, "Prompt · MCP · Skills", 1.0, 4.2, 11.0, 0.7,
                 font_size=28, color=COLOR_LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, "2026.03", 1.0, 5.5, 11.0, 0.5,
                 font_size=18, color=COLOR_LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # ── Slide 2: Overview ────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 0.08, 7.5, COLOR_HIGHLIGHT)
    add_rect(slide, 0.08, 0, 13.25, 1.4, COLOR_BG_ACCENT)
    add_text_box(slide, "综述", 0.25, 0.15, 12.5, 1.1,
                 font_size=36, bold=True, color=COLOR_HIGHLIGHT)

    add_text_box(
        slide,
        "目前 OT Flutter 客户端的 AI 提效主要使用三大工具，\n本文归纳了在使用这些工具过程中的提效心得。",
        0.5, 1.6, 12.0, 1.2, font_size=20, color=COLOR_LIGHT_GRAY
    )

    # Three blocks
    blocks = [
        ("01", "Prompt", "精准提示词 · 质疑设计\n及时 Review"),
        ("02", "MCP", "结构化知识输入\n连接组织文档与日志"),
        ("03", "Skills", "定义可复用工作流\n持续迭代改进"),
    ]
    colors = [RGBColor(0x0F, 0x3C, 0x78), RGBColor(0x0A, 0x4F, 0x6E), RGBColor(0x0A, 0x5F, 0x5A)]
    for i, (num, name, desc) in enumerate(blocks):
        x = 0.5 + i * 4.2
        add_rect(slide, x, 3.0, 3.8, 3.8, colors[i])
        add_text_box(slide, num, x + 0.2, 3.1, 1.0, 0.7,
                     font_size=28, bold=True, color=COLOR_HIGHLIGHT)
        add_text_box(slide, name, x + 0.2, 3.8, 3.4, 0.8,
                     font_size=28, bold=True, color=COLOR_WHITE)
        add_text_box(slide, desc, x + 0.2, 4.6, 3.4, 1.8,
                     font_size=16, color=COLOR_LIGHT_GRAY)

    # ── Slide 3: Prompt - 精准提示词 ────────────────────────────────────────
    add_bullet_slide(
        prs,
        "Prompt — 精准的提示词",
        [
            "目标：准确高效的设计，而非偶然的「惊喜」",
            "模糊提示词输出不稳定，需要反复修改",
            "清晰描述需求上下文、约束条件、期望输出格式",
            "将背景知识、已有代码结构一并提供给 AI",
        ]
    )

    # ── Slide 4: Prompt - 质疑 ───────────────────────────────────────────────
    add_two_column_slide(
        prs,
        "Prompt — 质疑与学习",
        left_header="询问为什么",
        left_items=[
            "不断学习进步的环节",
            "把握代码质量的重要方式",
            "理解 AI 的思考链路",
        ],
        right_header="质疑它的实现",
        right_items=[
            "发挥自身专业能力",
            "对设计方案提出合理质疑",
            "引导 AI 给出更优方案",
        ]
    )

    # ── Slide 5: Prompt - 代码案例 ───────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 0.08, 7.5, COLOR_HIGHLIGHT)
    add_rect(slide, 0.08, 0, 13.25, 1.4, COLOR_BG_ACCENT)
    add_text_box(slide, "Prompt — 案例：Enum 设计优化", 0.25, 0.15, 12.5, 1.1,
                 font_size=30, bold=True, color=COLOR_HIGHLIGHT)

    # Before label
    add_text_box(slide, "AI 初版（解耦但冗余）", 0.3, 1.55, 5.8, 0.45,
                 font_size=17, bold=True, color=RGBColor(0xFF, 0x80, 0x80))
    add_rect(slide, 0.3, 2.0, 5.8, 4.8, RGBColor(0x1E, 0x1E, 0x3A))
    before_code = (
        "enum VoiceClonePageSource {\n"
        "  homeNewUser,\n"
        "  homeReClone,\n"
        "  settings,\n"
        "  practice;\n"
        "}\n"
        "extension VoiceClonePageSourceX\n"
        "    : VoiceClonePageSource {\n"
        "  String get trackValue {\n"
        "    switch (this) {\n"
        "      case .homeNewUser: return '1';\n"
        "      ...\n"
        "    }\n"
        "  }\n"
        "}"
    )
    add_text_box(slide, before_code, 0.4, 2.05, 5.6, 4.7,
                 font_size=13, color=COLOR_LIGHT_GRAY)

    # Arrow
    add_text_box(slide, "质疑后 →", 6.3, 4.3, 1.0, 0.5,
                 font_size=18, bold=True, color=COLOR_YELLOW)

    # After label
    add_text_box(slide, "优化版（关联值，简洁）", 7.4, 1.55, 5.5, 0.45,
                 font_size=17, bold=True, color=COLOR_GREEN)
    add_rect(slide, 7.4, 2.0, 5.5, 3.5, RGBColor(0x1E, 0x1E, 0x3A))
    after_code = (
        "enum VoiceClonePageSource {\n"
        "  homeNewUser('1'),\n"
        "  homeReClone('2'),\n"
        "  settings('3'),\n"
        "  practice('4');\n\n"
        "  final String trackValue;\n"
        "  const VoiceClonePageSource(\n"
        "    this.trackValue);\n"
        "}"
    )
    add_text_box(slide, after_code, 7.5, 2.05, 5.3, 3.4,
                 font_size=13, color=COLOR_LIGHT_GRAY)
    add_text_box(slide, "新版 Dart 最佳实践：enum 直接关联埋点值",
                 7.4, 5.6, 5.5, 0.6, font_size=14, color=COLOR_YELLOW)

    # ── Slide 6: Prompt - 立刻 Review ─────────────────────────────────────────
    add_bullet_slide(
        prs,
        "Prompt — 立刻 Review 生成的代码",
        [
            "尽快发现问题，避免积累隐患",
            "防止 AI 意外修改无关代码",
            "两天后才发现 bug 排查成本成倍增加",
            "看不懂就「追问」，有疑虑就「质疑」",
            "完全确认后再提交，短期成本换长期效率",
        ]
    )

    # ── Slide 7: MCP Overview ────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 0.08, 7.5, COLOR_HIGHLIGHT)
    add_rect(slide, 0.08, 0, 13.25, 1.4, COLOR_BG_ACCENT)
    add_text_box(slide, "MCP — 协助 AI 访问组织知识", 0.25, 0.15, 12.5, 1.1,
                 font_size=30, bold=True, color=COLOR_HIGHLIGHT)

    mcps = [
        ("01", "埋点文档 MCP",
         "读取 Lark 埋点文档\n解析表格，结构化返回\n支持页面曝光/元素点击/结果返回",
         RGBColor(0x0F, 0x3C, 0x78)),
        ("02", "飞书文档转换\n(本项目)",
         "Lark ↔ Markdown 互转\n打通飞书文档与 Agent\nlark-doc-helper",
         RGBColor(0x0A, 0x4F, 0x6E)),
        ("03", "闪电日志 MCP",
         "AI 读取结构化日志\n按时间/等级/标签分段\n辅助问题分析排查",
         RGBColor(0x0A, 0x5F, 0x5A)),
    ]
    for i, (num, name, desc, color) in enumerate(mcps):
        x = 0.3 + i * 4.3
        add_rect(slide, x, 1.6, 4.0, 5.5, color)
        add_text_box(slide, num, x + 0.2, 1.7, 0.8, 0.6,
                     font_size=24, bold=True, color=COLOR_HIGHLIGHT)
        add_text_box(slide, name, x + 0.2, 2.3, 3.6, 0.9,
                     font_size=20, bold=True, color=COLOR_WHITE)
        add_text_box(slide, desc, x + 0.2, 3.3, 3.6, 2.8,
                     font_size=15, color=COLOR_LIGHT_GRAY)

    # ── Slide 8: MCP - 埋点结构化示例 ────────────────────────────────────────
    add_bullet_slide(
        prs,
        "MCP — 埋点文档结构化输出示例",
        subtitle="AI 直接读取结构化 JSON，免去人工整理埋点参数",
        bullets=[
            "分类：页面曝光 / 元素曝光 / 元素点击 / 结果返回",
            "字段：exclusive_id、$title、source 枚举值、business_type 等",
            "枚举值按版本标注迭代历史，AI 可感知历史版本演进",
            "AI 根据结构化信息自动生成符合规范的埋点代码",
        ]
    )

    # ── Slide 9: Skills Overview ─────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 0.08, 7.5, COLOR_HIGHLIGHT)
    add_rect(slide, 0.08, 0, 13.25, 1.4, COLOR_BG_ACCENT)
    add_text_box(slide, "Skills — 定义可复用工作流", 0.25, 0.15, 12.5, 1.1,
                 font_size=30, bold=True, color=COLOR_HIGHLIGHT)

    add_text_box(slide,
                 "把常用工作流定义为 Skills：有规则、有步骤的用 Prompt 描述，\n"
                 "结合 MCP / 内置 Scripts，精准控制特定问题的解决流程。",
                 0.5, 1.5, 12.0, 1.0, font_size=18, color=COLOR_LIGHT_GRAY)

    skills = [
        ("Figma → UI", "一键输入设计稿实现 UI\n5 个支付 A/B + 4 个挽回 A/B\n周级 → 1～2 天"),
        ("自动埋点", "每次迭代固定半天开销\n→ 几分钟完成\n准确率显著提升"),
    ]
    for i, (name, desc) in enumerate(skills):
        x = 1.0 + i * 6.0
        add_rect(slide, x, 2.7, 5.5, 4.4, RGBColor(0x0F, 0x3C, 0x78))
        add_text_box(slide, name, x + 0.3, 2.9, 4.9, 0.8,
                     font_size=26, bold=True, color=COLOR_YELLOW)
        add_text_box(slide, desc, x + 0.3, 3.8, 4.9, 2.5,
                     font_size=18, color=COLOR_WHITE)

    # ── Slide 10: 如何迭代 Skills ─────────────────────────────────────────────
    add_two_column_slide(
        prs,
        "Skills — 如何持续迭代",
        left_header="创建 Skill",
        left_items=[
            "使用元 Skill: create-skill",
            "从会话中提取步骤与思路",
            "/create-skill 总结这次会话...",
        ],
        right_header="迭代 Skill",
        right_items=[
            "完成开发后回顾遇到的问题",
            "将调整过程归纳进 Skill",
            "每次使用都让 Skill 更精准",
        ]
    )

    # ── Slide 11: Summary ────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 0.08, 7.5, COLOR_HIGHLIGHT)
    add_rect(slide, 0.08, 0, 13.25, 1.4, COLOR_BG_ACCENT)
    add_text_box(slide, "总结", 0.25, 0.15, 12.5, 1.1,
                 font_size=36, bold=True, color=COLOR_HIGHLIGHT)

    summary = [
        ("Prompt", "精准描述 + 及时质疑 + 立刻 Review = 稳定高质量输出"),
        ("MCP", "结构化知识注入 AI，消除文档和代码之间的信息鸿沟"),
        ("Skills", "固化最佳实践，让 AI 每次都「走正确的路」"),
    ]
    for i, (title, desc) in enumerate(summary):
        y = 1.8 + i * 1.6
        add_rect(slide, 0.3, y, 0.08, 1.0, COLOR_HIGHLIGHT)
        add_text_box(slide, title, 0.6, y, 2.5, 0.55,
                     font_size=24, bold=True, color=COLOR_YELLOW)
        add_text_box(slide, desc, 0.6, y + 0.55, 12.0, 0.8,
                     font_size=18, color=COLOR_WHITE)

    add_text_box(slide,
                 "三者结合，实现从「偶发的惊喜」到「稳定的提效」。",
                 0.5, 6.5, 12.0, 0.7, font_size=20,
                 bold=True, color=COLOR_HIGHLIGHT, align=PP_ALIGN.CENTER)

    out = "output.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build_ppt()
